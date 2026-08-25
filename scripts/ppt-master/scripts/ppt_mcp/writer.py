"""PPTX 写操作：文字替换、表格编辑、幻灯片复制、演讲者备注、转场、通用 plan 应用。

所有操作都走 `template_fill_pptx.apply_plan` 内部函数（不调 CLI），
绕过 `status==confirmed` 强制检查与时间戳后缀逻辑。

注意：apply_plan 按 plan.slides 逐条克隆页面并重建 sldIdLst，输出只含 plan
列出的页。因此任何"编辑已有整份"的操作都必须传全量页列表（未编辑页
pass-through），否则输出会丢失未列出的页面。

转场契约：_apply_plan 会把默认 effect 注入到所有未显式声明 transition 的页。
纯编辑操作（replace_text/replace_table_cells/duplicate_slide/add_notes）必须以
transition=None 调用 apply_plan 以保留源转场；set_transitions 仅改动目标页，
非目标页显式声明 "keep"。
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from template_fill_pptx import apply_plan as _apply_plan


def _full_deck_slides(pptx_path: str) -> list[int]:
    """返回 1..N 全部页号。"""
    prs = Presentation(str(pptx_path))
    return list(range(1, len(prs.slides) + 1))


def _new_plan() -> dict[str, Any]:
    return {"schema": "template_fill_pptx_fill_plan.v1", "status": "confirmed", "slides": []}


def apply_plan(
    pptxPath: str,
    plan: dict[str, Any],
    outputPath: str,
    transition: str | None = "fade",
    transitionDuration: float | None = 0.5,
) -> dict[str, Any]:
    """通用 plan 应用：直接调 template_fill_pptx.apply_plan 内部函数。

    MCP 端自动设置 status='confirmed'，无需二次确认。
    """
    src = Path(pptxPath).expanduser().resolve()
    if not src.exists():
        raise FileNotFoundError(f"PPTX not found: {src}")

    out = Path(outputPath).expanduser().resolve()
    out.parent.mkdir(parents=True, exist_ok=True)

    plan["status"] = "confirmed"

    # transition=None 时 _apply_plan 不为未显式声明转场的页注入默认 effect，
    # 从而保留源文件的转场；pptx_apply_plan 工具默认传 "fade"（见签名默认值）。
    _apply_plan(
        src,
        plan,
        out,
        transition=transition,
        transition_duration=transitionDuration if transitionDuration is not None else 0.5,
    )
    return {"outputPath": str(out), "slidesApplied": len(plan.get("slides", []))}


def replace_text(
    pptxPath: str,
    outputPath: str,
    replacements: list[dict[str, Any]],
    sourceSlide: int = 1,
    transition: str | None = None,
) -> dict[str, Any]:
    """按 slot_id / shape_id / shape_name 替换指定页的文字，其余页保留。"""
    plan = _new_plan()
    for idx in _full_deck_slides(pptxPath):
        slide: dict[str, Any] = {"source_slide": idx}
        if idx == sourceSlide:
            slide["replacements"] = replacements
            if transition:
                slide["transition"] = transition
        plan["slides"].append(slide)
    return apply_plan(pptxPath, plan, outputPath, transition=transition or None)


def replace_table_cells(
    pptxPath: str,
    outputPath: str,
    edits: list[dict[str, Any]],
    sourceSlide: int = 1,
    transition: str | None = None,
) -> dict[str, Any]:
    """替换指定页的表格单元格（edit 含 table_id/row/col/text），其余页保留。"""
    plan = _new_plan()
    for idx in _full_deck_slides(pptxPath):
        slide: dict[str, Any] = {"source_slide": idx}
        if idx == sourceSlide:
            slide["table_edits"] = edits
        if transition and idx == sourceSlide:
            slide["transition"] = transition
        plan["slides"].append(slide)
    return apply_plan(pptxPath, plan, outputPath, transition=transition or None)


def duplicate_slide(
    pptxPath: str,
    outputPath: str,
    slideIndex: int,
    count: int = 1,
    transition: str | None = None,
) -> dict[str, Any]:
    """复制指定页 N 次（插到末尾），其余页保留；count=0 时仅原样输出。"""
    plan = _new_plan()
    for idx in _full_deck_slides(pptxPath):
        plan["slides"].append({"source_slide": idx})
    for _ in range(max(0, count)):
        plan["slides"].append({"source_slide": slideIndex})
    return apply_plan(pptxPath, plan, outputPath, transition=transition or None)


def add_notes(
    pptxPath: str,
    outputPath: str,
    notes: list[dict[str, Any]],
    transition: str | None = None,
) -> dict[str, Any]:
    """给多个页面添加演讲者备注（notes: [{slideIndex, text}]），其余页保留。"""
    by_index = {item.get("slideIndex"): item.get("text", "") for item in notes}
    plan = _new_plan()
    for idx in _full_deck_slides(pptxPath):
        slide: dict[str, Any] = {"source_slide": idx}
        if idx in by_index:
            slide["notes"] = by_index[idx]
        plan["slides"].append(slide)
    return apply_plan(pptxPath, plan, outputPath, transition=transition or None)


def set_transitions(
    pptxPath: str,
    outputPath: str,
    transition: str,
    duration: float = 0.5,
    slides: list[int] | None = None,
) -> dict[str, Any]:
    """设置指定页（或全部页）的转场效果；未指定页保留源转场。"""
    targets = set(slides) if slides else set(_full_deck_slides(pptxPath))
    plan = _new_plan()
    for idx in _full_deck_slides(pptxPath):
        slide: dict[str, Any] = {"source_slide": idx}
        if idx in targets:
            slide["transition"] = transition
        else:
            # 显式 keep，防止默认 effect 覆盖未指定页的源转场
            slide["transition"] = "keep"
        plan["slides"].append(slide)
    return apply_plan(pptxPath, plan, outputPath, transition=None, transitionDuration=duration)