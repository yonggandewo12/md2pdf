"""PPTX 读取操作：结构概览、单页详情、文本提取、幻灯片转图片。

复用现有脚本库：
  - source_to_md.ppt_to_md.convert_presentation_to_markdown（文本提取）
  - pptx_to_svg.convert_pptx_to_svg（幻灯片 → SVG）
  - pymupdf（SVG → PNG 渲染，避免 cairosvg 的系统 libcairo 依赖）
"""

from __future__ import annotations

import contextlib
import io
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Emu


def _emu_to_inches(value: int | None) -> float | None:
    """EMU → 英寸（四舍五入到 2 位）。"""
    if value is None:
        return None
    return round(Emu(value).inches, 2)


def read_presentation(pptxPath: str) -> dict[str, Any]:
    """返回演示文稿总览：页数、尺寸、每页标题与 shape 数。"""
    path = Path(pptxPath).expanduser().resolve()
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {path}")

    prs = Presentation(str(path))
    slides = []
    for idx, slide in enumerate(prs.slides, 1):
        title = ""
        try:
            if slide.shapes.title is not None and slide.shapes.title.text:
                title = slide.shapes.title.text[:120]
        except Exception:
            pass
        slides.append({"index": idx, "title": title, "shapeCount": len(slide.shapes)})

    return {
        "slideCount": len(prs.slides),
        "slideWidthInches": _emu_to_inches(prs.slide_width),
        "slideHeightInches": _emu_to_inches(prs.slide_height),
        "slideWidthEmu": prs.slide_width,
        "slideHeightEmu": prs.slide_height,
        "slides": slides,
    }


def _shape_entry(shp: Any) -> dict[str, Any]:
    """把一个 shape 转成可序列化 dict。"""
    entry: dict[str, Any] = {
        "name": shp.name,
        "shapeId": getattr(shp, "shape_id", None),
        "shapeType": str(getattr(shp, "shape_type", "")),
        "left": _emu_to_inches(getattr(shp, "left", None)),
        "top": _emu_to_inches(getattr(shp, "top", None)),
        "width": _emu_to_inches(getattr(shp, "width", None)),
        "height": _emu_to_inches(getattr(shp, "height", None)),
    }
    try:
        if shp.has_text_frame:
            entry["text"] = shp.text_frame.text
    except Exception:
        pass
    try:
        if getattr(shp, "has_table", False):
            entry["table"] = [[cell.text for cell in row.cells] for row in shp.table.rows]
    except Exception:
        pass
    try:
        if getattr(shp, "shape_type", None) and "PICTURE" in str(getattr(shp, "shape_type", "")):
            entry["picture"] = True
    except Exception:
        pass
    return entry


def read_slide_details(pptxPath: str, slideIndex: int) -> dict[str, Any]:
    """返回指定页的 shapes 详情（名称/类型/位置/文本/表格）。"""
    path = Path(pptxPath).expanduser().resolve()
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {path}")

    prs = Presentation(str(path))
    count = len(prs.slides)
    if slideIndex < 1 or slideIndex > count:
        raise ValueError(f"slideIndex {slideIndex} out of range (1..{count})")

    slide = prs.slides[slideIndex - 1]
    shapes = [_shape_entry(shp) for shp in slide.shapes]
    return {"slideIndex": slideIndex, "slideCount": count, "shapes": shapes}


def extract_text(pptxPath: str, outputPath: str | None = None) -> dict[str, Any]:
    """把整份 PPTX 转成 Markdown 文本，可选写出文件。"""
    path = Path(pptxPath).expanduser().resolve()
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {path}")

    out = Path(outputPath).expanduser().resolve() if outputPath else path.with_suffix(".md")
    out.parent.mkdir(parents=True, exist_ok=True)

    from source_to_md.ppt_to_md import convert_presentation_to_markdown

    # 库函数内部向 stdout 打印 [INFO]/[OK] 日志，会污染单行 JSON 协议，
    # 重定向 stdout → stderr 保持协议干净。
    buf = io.StringIO()
    with contextlib.redirect_stdout(buf):
        markdown = convert_presentation_to_markdown(str(path), str(out))
    for line in buf.getvalue().splitlines():
        if line.strip():
            sys.stderr.write(line + "\n")

    return {"outputPath": str(out), "markdown": markdown[:50000]}


def to_images(
    pptxPath: str,
    outputDir: str | None = None,
    dpi: int = 150,
    format: str = "png",
    slides: str | None = None,
) -> dict[str, Any]:
    """每页幻灯片渲染为 PNG/JPEG 图片。

    流程：pptx_to_svg 生成自包含 SVG（inheritance_mode='flat'）→
    pymupdf 渲染为位图。targetPages 用与 PDF 相同的 "1-3,5,7-9" 语法。
    """
    path = Path(pptxPath).expanduser().resolve()
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {path}")

    out_dir = Path(outputDir).expanduser().resolve() if outputDir else path.with_name(f"{path.stem}_images")
    out_dir.mkdir(parents=True, exist_ok=True)

    from pptx_to_svg import convert_pptx_to_svg
    from pptx_to_svg.converter import ConvertOptions

    # flat 模式产出自包含 SVG（内联继承的母版/版式 shape），适合逐页渲染
    convert_pptx_to_svg(path, out_dir, ConvertOptions(inheritance_mode="flat"))

    svg_dir = out_dir / "svg-flat"
    if not svg_dir.exists():
        # 纯 flat 模式下 slide 直接写进 svg/（无 svg-flat/）
        svg_dir = out_dir / "svg"
    if not svg_dir.exists():
        raise RuntimeError(
            f"pptx_to_svg did not produce expected output directory (svg-flat/ or svg/) under {out_dir}"
        )

    import pymupdf

    ext = "jpg" if format == "jpeg" else "png"
    zoom = max(0.1, dpi) / 72.0
    matrix = pymupdf.Matrix(zoom, zoom)

    files: list[dict[str, Any]] = []
    svg_files = sorted(svg_dir.glob("slide_*.svg"), key=lambda p: int(p.stem.split("_")[1]))

    for svg_file in svg_files:
        idx = int(svg_file.stem.split("_")[1])
        if slides and idx not in _parse_targets(slides):
            continue
        doc = pymupdf.open(svg_file)
        try:
            page = doc[0]
            pix = page.get_pixmap(matrix=matrix)
            out_file = out_dir / f"slide_{idx:03d}.{ext}"
            pix.save(str(out_file))
            files.append({
                "path": str(out_file),
                "slideIndex": idx,
                "width": pix.width,
                "height": pix.height,
            })
        finally:
            doc.close()

    return {"outputDir": str(out_dir), "files": files}


def _parse_targets(spec: str) -> set[int]:
    """解析 "1-3,5,7-9" → {1,2,3,5,7,8,9}（1-based，越界不报错只跳过）。"""
    targets: set[int] = set()
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            lo_s, hi_s = part.split("-", 1)
            try:
                lo, hi = int(lo_s), int(hi_s)
            except ValueError:
                raise ValueError(f"Invalid slide range: '{part}' in '{spec}'")
            if lo > hi:
                raise ValueError(f"Descending range not allowed: '{part}' in '{spec}'")
            targets.update(range(lo, hi + 1))
        else:
            try:
                targets.add(int(part))
            except ValueError:
                raise ValueError(f"Invalid slide number: '{part}' in '{spec}'")
    return targets
